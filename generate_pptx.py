from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors
BLACK = RGBColor(0x0A, 0x0A, 0x0F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)   # violet
ACCENT2 = RGBColor(0x06, 0xB6, 0xD4)  # teal
MUTED = RGBColor(0x8B, 0x8B, 0xA0)
SURFACE = RGBColor(0x16, 0x16, 0x22)
GO = RGBColor(0x22, 0xC5, 0x5E)
MAYBE = RGBColor(0xF5, 0x9E, 0x0B)
SKIP = RGBColor(0xEF, 0x44, 0x44)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=BLACK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def accent_bar(slide, x, y, w=Inches(0.05), h=Inches(0.6)):
    box(slide, x, y, w, h, ACCENT)

def pill(slide, text, x, y, color=ACCENT):
    shape = slide.shapes.add_shape(5, x, y, Inches(1.6), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE

def dot(slide, x, y, color=ACCENT):
    shape = slide.shapes.add_shape(9, x, y, Inches(0.08), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ─────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────
s = add_slide()
bg(s)

# Glow circle
glow = s.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(6), Inches(6))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(0x2D, 0x10, 0x5A)
glow.line.fill.background()

box(s, Inches(0), Inches(6.9), W, Inches(0.6), SURFACE)

pill(s, "SALES INTELLIGENCE B2B", Inches(0.7), Inches(1.8))

txt(s, "REVORA", Inches(0.7), Inches(2.4), Inches(8), Inches(1.6), size=72, bold=True, color=WHITE)

# gradient text effect — subtitle
txt(s, "Transforme tes leads en briefs\ncommercaux prêts à l'emploi.", Inches(0.7), Inches(3.9), Inches(7), Inches(1.2), size=26, color=RGBColor(0xCC, 0xCC, 0xDD))

txt(s, "Scoring IA · Brief personnalisé · Export Excel", Inches(0.7), Inches(5.2), Inches(7), Inches(0.5), size=14, color=MUTED)

# Score badges on right
for i, (label, color, score) in enumerate([("GO", GO, "92"), ("MAYBE", MAYBE, "61"), ("SKIP", SKIP, "28")]):
    bx = Inches(9.8)
    by = Inches(1.8 + i * 1.5)
    b = s.shapes.add_shape(5, bx, by, Inches(2.8), Inches(1.1))
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = color
    b.line.width = Pt(1.5)
    tf = b.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = f"{label}  •  Score {score}"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = color

txt(s, "revora-sales.vercel.app", Inches(0.7), Inches(7.0), Inches(5), Inches(0.35), size=11, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 2 — PROBLÈME
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, ACCENT)

pill(s, "LE PROBLÈME", Inches(0.7), Inches(0.5), color=RGBColor(0x3B, 0x1A, 0x6A))
txt(s, "La qualification de leads\ncoûte trop cher en temps.", Inches(0.7), Inches(1.1), Inches(8), Inches(1.4), size=36, bold=True, color=WHITE)

problems = [
    ("⏱", "2 à 4h par semaine", "perdues à qualifier manuellement chaque lead"),
    ("📋", "Approches génériques", "le même template envoyé à tout le monde"),
    ("🎯", "Mauvaise priorisation", "les meilleurs leads contactés trop tard ou jamais"),
    ("❌", "Zéro contexte", "le commercial arrive en appel sans savoir à qui il parle"),
]

for i, (icon, title, desc) in enumerate(problems):
    bx = Inches(0.7 + (i % 2) * 6.1)
    by = Inches(2.9 + (i // 2) * 1.7)
    b = s.shapes.add_shape(5, bx, by, Inches(5.5), Inches(1.4))
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = RGBColor(0x2A, 0x2A, 0x3A)
    b.line.width = Pt(1)

    txt(s, icon, bx + Inches(0.2), by + Inches(0.15), Inches(0.5), Inches(0.5), size=20)
    txt(s, title, bx + Inches(0.7), by + Inches(0.12), Inches(4.5), Inches(0.4), size=14, bold=True, color=WHITE)
    txt(s, desc, bx + Inches(0.7), by + Inches(0.55), Inches(4.5), Inches(0.6), size=12, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 3 — SOLUTION
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, ACCENT2)

pill(s, "LA SOLUTION", Inches(0.7), Inches(0.5), color=RGBColor(0x03, 0x4A, 0x5A))
txt(s, "Un brief complet par lead,\ngénéré en quelques secondes.", Inches(0.7), Inches(1.1), Inches(9), Inches(1.4), size=36, bold=True, color=WHITE)

steps = [
    ("01", "Upload CSV", "Tu importes ton fichier de leads (nom, poste, entreprise, secteur, taille)"),
    ("02", "Définis ton ICP", "Tu décris ton client idéal : persona, problème, deal moyen, cycle de vente"),
    ("03", "L'IA analyse", "Gemini 2.5 Flash croise chaque lead avec ton ICP et génère le brief"),
    ("04", "Export & action", "Tu exportes en Excel mis en forme et tu commences à prospecter"),
]

for i, (num, title, desc) in enumerate(steps):
    bx = Inches(0.7 + i * 3.1)
    by = Inches(3.0)

    # Connector line
    if i < 3:
        line = s.shapes.add_shape(1, bx + Inches(2.6), by + Inches(0.4), Inches(0.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

    b = s.shapes.add_shape(5, bx, by, Inches(2.7), Inches(3.2))
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = ACCENT
    b.line.width = Pt(1)

    txt(s, num, bx + Inches(0.2), by + Inches(0.2), Inches(0.6), Inches(0.5), size=11, bold=True, color=ACCENT)
    txt(s, title, bx + Inches(0.2), by + Inches(0.65), Inches(2.3), Inches(0.5), size=15, bold=True, color=WHITE)
    txt(s, desc, bx + Inches(0.2), by + Inches(1.15), Inches(2.3), Inches(1.8), size=11, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 4 — CE QUE TU REÇOIS
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, ACCENT)

pill(s, "LE BRIEF", Inches(0.7), Inches(0.5))
txt(s, "Ce que REVORA génère\npour chaque lead.", Inches(0.7), Inches(1.1), Inches(6), Inches(1.4), size=36, bold=True, color=WHITE)

items = [
    ("🎯", "Score 0-100 + Priorité", "GO / MAYBE / SKIP / VETO avec justification"),
    ("📝", "Briefing", "Pourquoi ce lead matche (ou pas) ton ICP"),
    ("💬", "Première phrase d'ouverture", "Rédigée pour ce contact précis, pas un template"),
    ("📡", "Canal recommandé", "Cold Call, LinkedIn, Email ou Multi-touch + raison"),
    ("🛡", "3 objections anticipées", "Avec la réponse adaptée à chaque objection"),
    ("⏰", "Timing idéal", "Quand contacter ce profil"),
    ("⚠️", "Piège à éviter", "L'erreur classique sur ce type de prospect"),
]

for i, (icon, title, desc) in enumerate(items):
    col = i // 4
    row = i % 4
    bx = Inches(0.7 + col * 6.4)
    by = Inches(2.8 + row * 1.1)

    dot(s, bx, by + Inches(0.15), ACCENT)
    txt(s, title, bx + Inches(0.25), by, Inches(5.5), Inches(0.4), size=13, bold=True, color=WHITE)
    txt(s, desc, bx + Inches(0.25), by + Inches(0.38), Inches(5.5), Inches(0.5), size=11, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 5 — SCORING EXPLIQUÉ
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, MAYBE)

pill(s, "LE SCORING", Inches(0.7), Inches(0.5), color=RGBColor(0x5A, 0x3A, 0x03))
txt(s, "Comment le score\nest calculé.", Inches(0.7), Inches(1.1), Inches(7), Inches(1.4), size=36, bold=True, color=WHITE)

txt(s, "Le modèle reçoit les données brutes du lead + ton ICP et croise :", Inches(0.7), Inches(2.7), Inches(11), Inches(0.4), size=14, color=MUTED)

criteria = [
    "Le titre de poste correspond-il au décideur ou influenceur que tu cibles ?",
    "La taille d'entreprise entre-t-elle dans ta cible ?",
    "Le secteur est-il dans ta roue ?",
    "Cette personne a-t-elle probablement le problème que tu résous ?",
    "Y a-t-il un signal bloquant (concurrent, trop petit, hors budget) ?",
]

for i, c in enumerate(criteria):
    by = Inches(3.2 + i * 0.65)
    dot(s, Inches(0.7), by + Inches(0.12), MAYBE)
    txt(s, c, Inches(1.1), by, Inches(8), Inches(0.55), size=13, color=WHITE)

# Score legend
for i, (label, color, range_) in enumerate([("GO", GO, "75 – 100"), ("MAYBE", MAYBE, "40 – 74"), ("SKIP", SKIP, "0 – 39"), ("VETO", RGBColor(0x6B, 0x21, 0xA8), "Critère bloquant")]):
    bx = Inches(9.5)
    by = Inches(2.6 + i * 1.1)
    b = s.shapes.add_shape(5, bx, by, Inches(3.0), Inches(0.9))
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = color
    b.line.width = Pt(2)
    txt(s, label, bx + Inches(0.2), by + Inches(0.05), Inches(1.2), Inches(0.4), size=14, bold=True, color=color)
    txt(s, range_, bx + Inches(0.2), by + Inches(0.45), Inches(2.5), Inches(0.35), size=11, color=MUTED)

txt(s, "⚠  La qualité du scoring = la qualité de ton CSV + la précision de ton ICP", Inches(0.7), Inches(6.6), Inches(11), Inches(0.5), size=11, color=MUTED, italic=True)

# ─────────────────────────────────────────────
# SLIDE 6 — PERSONALISATION
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, ACCENT)

pill(s, "PERSONNALISATION", Inches(0.7), Inches(0.5))
txt(s, "Pourquoi chaque brief\nest vraiment unique.", Inches(0.7), Inches(1.1), Inches(8), Inches(1.4), size=36, bold=True, color=WHITE)

# Left: what's fed in
box(s, Inches(0.7), Inches(2.7), Inches(5.5), Inches(4.0), SURFACE)
txt(s, "CE QUE LE MODÈLE REÇOIT", Inches(0.85), Inches(2.85), Inches(5), Inches(0.4), size=10, bold=True, color=ACCENT)

example = [
    ('Poste :', 'Head of Sales'),
    ('Entreprise :', 'Doctrine (LegalTech)'),
    ('Taille :', '50-200 salariés'),
    ('ICP :', 'Scale-ups avec équipe SDR sans process de qualification'),
    ('Deal moyen :', '15 000 €  |  Cycle : 3 mois'),
]
for i, (k, v) in enumerate(example):
    by = Inches(3.3 + i * 0.55)
    txt(s, k, Inches(0.9), by, Inches(1.5), Inches(0.45), size=11, color=MUTED)
    txt(s, v, Inches(2.2), by, Inches(3.8), Inches(0.45), size=11, bold=True, color=WHITE)

# Arrow
txt(s, "→", Inches(6.4), Inches(4.3), Inches(0.5), Inches(0.5), size=28, color=ACCENT, align=PP_ALIGN.CENTER)

# Right: output example
box(s, Inches(7.0), Inches(2.7), Inches(5.6), Inches(4.0), SURFACE)
txt(s, "OUVERTURE GÉNÉRÉE", Inches(7.15), Inches(2.85), Inches(5), Inches(0.4), size=10, bold=True, color=ACCENT)
txt(s, '"Thomas, j\'ai vu que Doctrine a levé en série B l\'an dernier — dans ce contexte de scale, structurer le process de qualification de vos SDR avant d\'augmenter l\'équipe peut faire la différence entre 6 et 12 mois pour atteindre le quota. C\'est exactement ce qu\'on résout."',
    Inches(7.15), Inches(3.35), Inches(5.2), Inches(2.0), size=11, color=WHITE, italic=True)

txt(s, "Canal : LinkedIn  |  Timing : Post-levée, avant recrutement SDR", Inches(7.15), Inches(5.5), Inches(5.2), Inches(0.4), size=10, color=MAYBE)

txt(s, "Pas un template. Pas générique. Rédigé à partir des données réelles du lead.", Inches(0.7), Inches(7.0), Inches(11), Inches(0.35), size=11, color=MUTED, italic=True)

# ─────────────────────────────────────────────
# SLIDE 7 — POUR QUI
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, ACCENT2)

pill(s, "POUR QUI ?", Inches(0.7), Inches(0.5), color=RGBColor(0x03, 0x4A, 0x5A))
txt(s, "REVORA est fait pour\nles équipes commerciales B2B.", Inches(0.7), Inches(1.1), Inches(9), Inches(1.4), size=36, bold=True, color=WHITE)

profiles = [
    ("SDR / BDR", "Tu gères des listes de leads et tu perds du temps à qualifier avant de prospecter."),
    ("Account Executive", "Tu veux arriver en appel avec du contexte, pas juste un nom et une boîte."),
    ("Head of Sales", "Tu veux que ton équipe passe plus de temps à vendre et moins à préparer."),
    ("Fondateur / Solo", "Tu fais tout toi-même — REVORA te donne le niveau de préparation d'une équipe."),
]

for i, (title, desc) in enumerate(profiles):
    col = i % 2
    row = i // 2
    bx = Inches(0.7 + col * 6.3)
    by = Inches(3.0 + row * 1.8)
    b = s.shapes.add_shape(5, bx, by, Inches(5.8), Inches(1.5))
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = ACCENT2
    b.line.width = Pt(1)
    txt(s, title, bx + Inches(0.25), by + Inches(0.2), Inches(5.3), Inches(0.45), size=15, bold=True, color=WHITE)
    txt(s, desc, bx + Inches(0.25), by + Inches(0.7), Inches(5.3), Inches(0.65), size=12, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 8 — TARIFS
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, ACCENT)

pill(s, "TARIFS", Inches(0.7), Inches(0.5))
txt(s, "Simple. Sans surprise.", Inches(0.7), Inches(1.1), Inches(8), Inches(0.9), size=36, bold=True, color=WHITE)

plans = [
    ("Solo", "29€", "/mois", "1 utilisateur · 500 leads/mois · Export Excel · Support email", False),
    ("Pro", "79€", "/mois", "5 utilisateurs · 3 000 leads/mois · ICP sauvegardé · Support prioritaire · 14j gratuits", True),
    ("Scale", "149€", "/mois", "Utilisateurs illimités · Leads illimités · API · SLA 99,9% · Account manager", False),
]

for i, (name, price, per, features, highlight) in enumerate(plans):
    bx = Inches(0.7 + i * 4.2)
    by = Inches(2.2)
    bw = Inches(3.8)
    bh = Inches(4.6)
    b = s.shapes.add_shape(5, bx, by, bw, bh)
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = ACCENT if highlight else RGBColor(0x2A, 0x2A, 0x3A)
    b.line.width = Pt(2 if highlight else 1)

    if highlight:
        pill_b = s.shapes.add_shape(5, bx + Inches(0.8), by - Inches(0.2), Inches(2.2), Inches(0.38))
        pill_b.fill.solid()
        pill_b.fill.fore_color.rgb = ACCENT
        pill_b.line.fill.background()
        tf = pill_b.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = "LE PLUS POPULAIRE"
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = WHITE

    txt(s, name, bx + Inches(0.25), by + Inches(0.3), bw - Inches(0.4), Inches(0.45), size=16, bold=True, color=WHITE)
    txt(s, price, bx + Inches(0.25), by + Inches(0.85), Inches(1.5), Inches(0.7), size=36, bold=True, color=ACCENT if highlight else WHITE)
    txt(s, per, bx + Inches(1.6), by + Inches(1.25), Inches(1.8), Inches(0.35), size=12, color=MUTED)

    for j, feat in enumerate(features.split(" · ")):
        fy = by + Inches(1.85 + j * 0.52)
        dot(s, bx + Inches(0.25), fy + Inches(0.12), ACCENT if highlight else ACCENT2)
        txt(s, feat, bx + Inches(0.5), fy, bw - Inches(0.6), Inches(0.45), size=11, color=WHITE if highlight else MUTED)

txt(s, "−20 % sur tous les plans en facturation annuelle", Inches(0.7), Inches(7.05), Inches(11), Inches(0.35), size=11, color=MUTED, italic=True)

# ─────────────────────────────────────────────
# SLIDE 9 — CONFIANCE / SÉCURITÉ
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
box(s, Inches(0), Inches(0), Inches(0.06), H, ACCENT2)

pill(s, "CONFIANCE & SÉCURITÉ", Inches(0.7), Inches(0.5), color=RGBColor(0x03, 0x4A, 0x5A))
txt(s, "Tes données sont les tiennes,\npas celles de l'IA.", Inches(0.7), Inches(1.1), Inches(9), Inches(1.4), size=36, bold=True, color=WHITE)

guarantees = [
    ("🔒", "Données isolées", "Chaque utilisateur accède uniquement à ses propres données. Zéro partage entre comptes."),
    ("🚫", "Pas d'entraînement IA", "Les appels à Gemini se font via un abonnement pro qui exclut l'utilisation des données pour entraîner les modèles."),
    ("🇪🇺", "Conforme RGPD", "Minimisation des données, chiffrement TLS, droit d'accès / suppression sur demande à privacy@revora.app."),
    ("👁", "Scoring transparent", "Pour chaque lead, tu vois le raisonnement complet. Pas de boîte noire."),
]

for i, (icon, title, desc) in enumerate(guarantees):
    col = i % 2
    row = i // 2
    bx = Inches(0.7 + col * 6.3)
    by = Inches(2.9 + row * 1.9)
    txt(s, icon, bx, by, Inches(0.5), Inches(0.5), size=22)
    txt(s, title, bx + Inches(0.6), by + Inches(0.05), Inches(5.4), Inches(0.4), size=14, bold=True, color=WHITE)
    txt(s, desc, bx + Inches(0.6), by + Inches(0.5), Inches(5.4), Inches(1.1), size=12, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 10 — CTA / CLOSING
# ─────────────────────────────────────────────
s = add_slide()
bg(s)

glow2 = s.shapes.add_shape(9, Inches(3), Inches(1), Inches(8), Inches(8))
glow2.fill.solid()
glow2.fill.fore_color.rgb = RGBColor(0x1A, 0x08, 0x35)
glow2.line.fill.background()

txt(s, "REVORA", Inches(1), Inches(1.2), Inches(11), Inches(1.2), size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Commence à prospecter comme un senior.", Inches(1), Inches(2.6), Inches(11), Inches(0.7), size=24, color=RGBColor(0xCC, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

txt(s, "14 jours gratuits · Sans carte bancaire · Accès sous 24h", Inches(1), Inches(3.5), Inches(11), Inches(0.5), size=14, color=MUTED, align=PP_ALIGN.CENTER)

# CTA button
btn = s.shapes.add_shape(5, Inches(4.4), Inches(4.4), Inches(4.5), Inches(0.75))
btn.fill.solid()
btn.fill.fore_color.rgb = ACCENT
btn.line.fill.background()
tf = btn.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = "revora-sales.vercel.app"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = WHITE

txt(s, "Questions ?  →  support@revora.app", Inches(1), Inches(5.6), Inches(11), Inches(0.4), size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Stats
for i, (val, label) in enumerate([("< 5 min", "pour lancer\nune analyse"), ("200 leads", "par analyse\nmax"), ("7 données", "générées\npar lead")]):
    bx = Inches(1.5 + i * 3.8)
    by = Inches(6.2)
    txt(s, val, bx, by, Inches(3.2), Inches(0.5), size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, label, bx, by + Inches(0.5), Inches(3.2), Inches(0.6), size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Save
output = r"C:\Users\taffi\Desktop\REVORA_Presentation.pptx"
prs.save(output)
print(f"OK - Presentation saved: {output}")
